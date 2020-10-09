# encoding: utf-8
import fitz
import time
from tqdm import tqdm
import shutil
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
import sys
# import zipfile


def pdfToPPF(filepath):
    if bool(filepath) != True:
        return False
    # 根据路径获取获取pdf名称
    filename = os.path.basename(filepath)
    # 将pdf一张张切割成jpg
    timestamp = time.time()
    print('切割pdf为jpg...')
    # 打开一个PDF文件，doc为Document类型，是一个包含每一页PDF文件的列表
    doc = fitz.open(filepath)
    # 设置图片的旋转角度
    rotate = int(0)
    # 设置图片相对于PDF文件在X轴上的缩放比例
    zoom_x = 2.0
    # 设置图片相对于PDF文件在Y轴上的缩放比例
    zoom_y = 2.0
    trans = fitz.Matrix(zoom_x, zoom_y).preRotate(rotate)

    # 后缀字符串
    suffix = ".pdf"
    filenameLen = len(filename)
    # pdf大小写逻辑处理
    if filename.endswith(".PDF", filenameLen - 4, filenameLen) == True:
        suffix = ".PDF"

    # 保证输出的文件名不变 且不重复
    new_full_name = filename.replace(suffix, "", 1) + str(round(timestamp * 1000))

    # 获取根路径
    base_path = getBast_path()
    print(base_path)
    # 图片路径
    jpg_path = getJpgPath(new_full_name)
    print(jpg_path)
    print("%s开始转换..." % filename)
    if doc.pageCount > 1:  # 获取PDF的页数
        for pg in tqdm(range(doc.pageCount)):
            page = doc[pg]  # 获得第pg页
            pm = page.getPixmap(matrix=trans, alpha=False)  # 将其转化为光栅文件（位数）
            if not os.path.exists(jpg_path):
                os.mkdir(jpg_path)  # makedirs
            pm.writeImage(os.path.join(jpg_path, '%s-%s.jpg' % (new_full_name, pg)))  # 将其输入为相应的图片格式，可以为位图，也可以为矢量图
            # 我本来想输出为jpg文件，但是在网页中都是png格式（即调用writePNG），再转换成别的图像文件前，最好查一下是否支持
    else:
        page = doc[0]
        pm = page.getPixmap(matrix=trans, alpha=False)
        pm.writeImage("%s.jpg" % new_full_name)
    print("%s转换jpg完成！" % filename)
    print('耗时：', time.time() - timestamp, 's')
    print('\n\n')

    # ================================================================================================================================================
    # 将分割好的jpg图片整合到ppt
    timestamp = time.time()
    print('整合jpg为ppt...')

    pages = os.listdir(jpg_path)
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    for index, page in enumerate(tqdm(pages)):
        # print(index)
        # Save as 'jpg' in jpgs dir
        jpg_file = os.path.join(jpg_path, '%s-%d.jpg' % (new_full_name, index))
        # Get width/height of image
        image = Image.open(jpg_file)
        height = image.height
        width = image.width

        # #Rotate 270 degrees if horizontal
        # if height > width:
        #     adjusted = image.rotate(270, expand=True)
        #     adjusted.save(jpg_file)

        # Setup slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        left = top = 0
        slide.shapes.add_picture(jpg_file, left, top, height=prs.slide_height, width=prs.slide_width)
        # 必须退出  不然删除会异常
        image.close()

    new_full_path = os.path.join(base_path, '__result', '%s.pptx' % new_full_name)
    prs.save(new_full_path)

    print("成功保存ppt文件 %s.pptx" % new_full_name)
    print('耗时：', time.time() - timestamp)

    # 文件压缩
    # file_zip_path = zip_file_path(new_full_path)

    # 删除缓存的图片文件
    delBypath(jpg_path)
    # 删除未压缩的ppt文件
    # os.remove(new_full_path)
    print("finish:" + new_full_path)
    # 退出进程
    sys.exit(0)


# 文件压缩
def zip_file_path(file_path):
    """
    压缩文件，相对路径
    :param filepath: 压缩的文件夹路径
    :return:
    """
    # # 获取文件名
    # file_name = os.path.basename(file_path)
    # # 文件路径
    # file_base_path = os.path.split(file_path)[0]
    #
    # filelists = os.listdir(file_base_path)
    #
    # f = zipfile.ZipFile(file_path + '.zip', 'w', compression = zipfile.ZIP_DEFLATED)
    #
    # for fil in filelists:
    #     filefullpath = os.path.join(file_base_path, fil)
    #     # filefullpath是文件的全路径，fil是文件名，这样就不会带目录啦
    #     f.write(filefullpath, fil)
    # # 调用了close方法才会保证完成压缩
    # f.close()
    # return file_path + '.zip'


# 根路径获取
def getBast_path():
    bast_path = os.path.split(os.path.abspath('.'))[1]
    if bast_path == 'PDFtoPPT':
        path = os.path.split(os.path.abspath('.'))[0]
    else:
        path = os.path.abspath('.')
    return os.path.join(path, 'upload')


# 获取文件路径
def getJpgPath(str_path):
    return os.path.join(getBast_path(), '__jpgs', '%s' % str_path)


# 文件夹删除操作
def delBypath(filePath):
    shutil.rmtree(filePath)
    return '0'


# print(r'' + sys.argv[1])
pdfToPPF(os.path.join(r'' + sys.argv[1]))
